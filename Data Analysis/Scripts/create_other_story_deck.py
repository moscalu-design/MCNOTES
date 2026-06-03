from __future__ import annotations

import json
import math
from collections import Counter
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DATA_PATH = ROOT / "Data Analysis" / "dashboard_data.json"
ASSET_DIR = ROOT / "Data Analysis" / "Exports" / "other_story_deck_assets"
OUTPUT_PATH = ROOT / "OTHER_MC_Notes_Story.pptx"

DEPARTMENTS = [
    "OPS",
    "GLO",
    "PJ",
    "RM",
    "OCCO",
    "JU",
    "ECON",
    "CFC",
    "EIF",
    "FI",
    "IG",
    "PMM",
    "SG",
    "GIS",
    "HR",
    "OTHER",
]

MC_NOTE_TYPES = ["NOTEMCDEC", "NOTEMCDISC", "NOTEMCINFO"]
MC_LABELS = {
    "NOTEMCDEC": "Decision",
    "NOTEMCDISC": "Discussion",
    "NOTEMCINFO": "Info",
}

COLORS = {
    "blue": "#376996",
    "teal": "#2a9d8f",
    "red": "#c1666b",
    "yellow": "#b38b00",
    "ink": "#17212b",
    "muted": "#607080",
    "grid": "#dce3eb",
    "paper": "#ffffff",
    "bg": "#f3f6fa",
}
TYPE_COLORS = {
    "NOTEMCDEC": COLORS["blue"],
    "NOTEMCDISC": COLORS["red"],
    "NOTEMCINFO": COLORS["teal"],
}
PAGE_COLORS = {
    "Pre-opinion": COLORS["blue"],
    "Opinion": COLORS["teal"],
    "Annex": COLORS["yellow"],
}


def n(value) -> float:
    try:
        out = float(value)
    except (TypeError, ValueError):
        return 0.0
    return out if math.isfinite(out) else 0.0


def fmt_int(value: float) -> str:
    return f"{value:,.0f}"


def fmt_one(value: float) -> str:
    return f"{value:,.1f}"


def fmt_pct(value: float) -> str:
    return f"{value:,.1f}%"


def clean_axis(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#ccd5df")
    ax.spines["bottom"].set_color("#ccd5df")
    ax.tick_params(colors=COLORS["muted"], labelsize=9)
    ax.grid(True, axis="y", color=COLORS["grid"], linewidth=0.8)
    ax.set_axisbelow(True)


def save_fig(fig, name: str) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / f"{name}.png"
    fig.savefig(path, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def load_other_rows() -> pd.DataFrame:
    payload = json.loads(DATA_PATH.read_text(encoding="utf-8"))
    df = pd.DataFrame(payload["records"])
    df = df[df["Template Type"].eq("OTHER")].copy()
    df["BO Validation Date"] = pd.to_datetime(df["BO Validation Date"], errors="coerce")
    df["BO Validation Month"] = df["BO Validation Date"].dt.to_period("M").astype("string")
    df["BO Validation Quarter"] = df["BO Validation Date"].dt.to_period("Q").astype("string")
    for col in ["Document Page Count", "Page count before opinion", "Annex Page Count", "Text Before Opinions", *DEPARTMENTS]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def has_service_opinion(row: pd.Series) -> bool:
    return any(n(row.get(service)) > 0 for service in DEPARTMENTS)


def service_opinion_frame(df: pd.DataFrame) -> pd.DataFrame:
    rows = []
    for _, row in df.iterrows():
        for service in DEPARTMENTS:
            words = n(row.get(service))
            if words > 0:
                rows.append(
                    {
                        "service": service,
                        "words": words,
                        "file": row.get("File Name", ""),
                        "month": row.get("BO Validation Month"),
                    }
                )
    return pd.DataFrame(rows)


def build_service_volume_chart(df: pd.DataFrame) -> Path:
    opinions = service_opinion_frame(df)
    stats = (
        opinions.groupby("service")
        .agg(opinions=("words", "size"), mean_words=("words", "mean"), median_words=("words", "median"))
        .sort_values(["opinions", "median_words"], ascending=[False, False])
        .head(12)
        .reset_index()
    )
    stats = stats.iloc[::-1]

    fig, (ax_count, ax_words) = plt.subplots(
        1,
        2,
        figsize=(12.2, 5.0),
        gridspec_kw={"width_ratios": [1.45, 1.0], "wspace": 0.12},
    )
    y = np.arange(len(stats))
    ax_count.barh(y, stats["opinions"], color=COLORS["blue"], alpha=0.58, height=0.62)
    ax_count.set_yticks(y, stats["service"], fontweight="bold", color=COLORS["ink"])
    ax_count.set_xlabel("Opinion rows", color=COLORS["muted"], fontsize=9)
    ax_count.set_title("Volume", loc="left", color=COLORS["ink"], fontsize=12, fontweight="bold")
    ax_count.grid(True, axis="x", color=COLORS["grid"], linewidth=0.8)
    ax_count.grid(False, axis="y")
    for i, value in enumerate(stats["opinions"]):
        ax_count.text(value + max(stats["opinions"]) * 0.015, i, fmt_int(value), va="center", fontsize=8.5, color=COLORS["muted"], fontweight="bold")
    ax_count.spines[["top", "right", "left"]].set_visible(False)
    ax_count.spines["bottom"].set_color("#ccd5df")
    ax_count.tick_params(axis="x", colors=COLORS["muted"], labelsize=8.5)
    ax_count.tick_params(axis="y", length=0)

    ax_words.hlines(y, stats["median_words"], stats["mean_words"], color="#aeb8c4", linewidth=2.2)
    ax_words.scatter(stats["mean_words"], y, s=52, color=COLORS["ink"], label="Mean", zorder=3)
    ax_words.scatter(stats["median_words"], y, s=52, color="white", edgecolor=COLORS["ink"], linewidth=1.7, label="Median", zorder=4)
    ax_words.set_yticks(y, [""] * len(y))
    ax_words.set_xlabel("Words per opinion", color=COLORS["muted"], fontsize=9)
    ax_words.set_title("Length", loc="left", color=COLORS["ink"], fontsize=12, fontweight="bold")
    ax_words.grid(True, axis="x", color=COLORS["grid"], linewidth=0.8)
    ax_words.grid(False, axis="y")
    ax_words.spines[["top", "right", "left"]].set_visible(False)
    ax_words.spines["bottom"].set_color("#ccd5df")
    ax_words.tick_params(axis="x", colors=COLORS["muted"], labelsize=8.5)
    ax_words.tick_params(axis="y", length=0)
    ax_words.legend(loc="lower right", frameon=False, fontsize=9)

    return save_fig(fig, "other_service_volume_length")


def page_composition_monthly(df: pd.DataFrame) -> pd.DataFrame:
    dated = df.dropna(subset=["BO Validation Date"]).copy()
    grouped = (
        dated.groupby("BO Validation Month")
        .agg(
            documents=("File Name", "count"),
            total_pages=("Document Page Count", "mean"),
            pre_opinion=("Page count before opinion", "mean"),
            annex=("Annex Page Count", "mean"),
        )
        .reset_index()
        .sort_values("BO Validation Month")
    )
    grouped["opinion"] = (grouped["total_pages"] - grouped["pre_opinion"] - grouped["annex"]).clip(lower=0)
    return grouped


def build_page_composition_chart(df: pd.DataFrame) -> Path:
    grouped = page_composition_monthly(df)
    x = np.arange(len(grouped))
    labels = grouped["BO Validation Month"].astype(str).tolist()
    tick_step = max(1, math.ceil(len(labels) / 12))

    fig, (ax_stack, ax_ma) = plt.subplots(2, 1, figsize=(12.2, 6.2), gridspec_kw={"height_ratios": [2.2, 1.0], "hspace": 0.32})
    bottom = np.zeros(len(grouped))
    parts = [
        ("Pre-opinion", grouped["pre_opinion"].to_numpy()),
        ("Opinion", grouped["opinion"].to_numpy()),
        ("Annex", grouped["annex"].to_numpy()),
    ]
    for label, values in parts:
        ax_stack.bar(x, values, bottom=bottom, color=PAGE_COLORS[label], label=label, width=0.72)
        bottom += values
    ax_stack.set_title("Monthly mean page composition", loc="left", color=COLORS["ink"], fontsize=12, fontweight="bold")
    ax_stack.set_ylabel("Mean pages/document", color=COLORS["muted"], fontsize=9)
    ax_stack.set_xticks(x[::tick_step], labels[::tick_step], rotation=0)
    clean_axis(ax_stack)
    ax_stack.legend(ncols=3, loc="upper right", frameon=False, fontsize=9)

    for label, values in parts:
        rolling = pd.Series(values).rolling(3, min_periods=1).mean()
        ax_ma.plot(x, rolling, color=PAGE_COLORS[label], linewidth=2.4, label=label)
    ax_ma.set_title("3-month moving average", loc="left", color=COLORS["ink"], fontsize=11, fontweight="bold")
    ax_ma.set_ylabel("Mean pages", color=COLORS["muted"], fontsize=9)
    ax_ma.set_xticks(x[::tick_step], labels[::tick_step], rotation=0)
    clean_axis(ax_ma)

    return save_fig(fig, "other_monthly_page_composition")


def valid_mc_df(df: pd.DataFrame) -> pd.DataFrame:
    return df[df["MC_Note_Type"].isin(MC_NOTE_TYPES)].copy()


def build_mc_monthly_chart(df: pd.DataFrame) -> Path:
    dated = valid_mc_df(df).dropna(subset=["BO Validation Date"])
    counts = (
        dated.groupby(["BO Validation Month", "MC_Note_Type"])
        .size()
        .unstack(fill_value=0)
        .reindex(columns=MC_NOTE_TYPES, fill_value=0)
        .sort_index()
    )
    x = np.arange(len(counts))
    labels = counts.index.astype(str).tolist()
    tick_step = max(1, math.ceil(len(labels) / 12))

    fig, (ax_stack, ax_ma) = plt.subplots(2, 1, figsize=(12.2, 6.0), gridspec_kw={"height_ratios": [2.1, 1.0], "hspace": 0.33})
    bottom = np.zeros(len(counts))
    for note_type in MC_NOTE_TYPES:
        values = counts[note_type].to_numpy()
        ax_stack.bar(x, values, bottom=bottom, width=0.72, color=TYPE_COLORS[note_type], label=MC_LABELS[note_type])
        bottom += values
    ax_stack.set_title("Monthly documents by MC note type", loc="left", color=COLORS["ink"], fontsize=12, fontweight="bold")
    ax_stack.set_ylabel("Documents", color=COLORS["muted"], fontsize=9)
    ax_stack.set_xticks(x[::tick_step], labels[::tick_step])
    clean_axis(ax_stack)
    ax_stack.legend(ncols=3, loc="upper right", frameon=False, fontsize=9)

    for note_type in MC_NOTE_TYPES:
        rolling = counts[note_type].rolling(3, min_periods=1).mean()
        ax_ma.plot(x, rolling, color=TYPE_COLORS[note_type], linewidth=2.4, label=MC_LABELS[note_type])
    ax_ma.set_title("3-month moving average", loc="left", color=COLORS["ink"], fontsize=11, fontweight="bold")
    ax_ma.set_ylabel("Documents", color=COLORS["muted"], fontsize=9)
    ax_ma.set_xticks(x[::tick_step], labels[::tick_step])
    clean_axis(ax_ma)

    return save_fig(fig, "other_monthly_mc_note_type")


def build_quarterly_opinion_chart(df: pd.DataFrame) -> Path:
    dated = valid_mc_df(df).dropna(subset=["BO Validation Date"]).copy()
    dated["Has Service Opinion"] = dated.apply(has_service_opinion, axis=1)
    grouped = (
        dated.groupby(["BO Validation Quarter", "MC_Note_Type"])
        .agg(documents=("File Name", "count"), opinion_docs=("Has Service Opinion", "sum"))
        .reset_index()
    )
    grouped["share"] = grouped["opinion_docs"] / grouped["documents"] * 100
    quarters = sorted(grouped["BO Validation Quarter"].astype(str).unique())
    x = np.arange(len(quarters))

    fig, ax = plt.subplots(figsize=(12.2, 4.8))
    for note_type in MC_NOTE_TYPES:
        series = (
            grouped[grouped["MC_Note_Type"].eq(note_type)]
            .set_index("BO Validation Quarter")["share"]
            .reindex(quarters)
        )
        ax.plot(x, series, color=TYPE_COLORS[note_type], linewidth=2.8, marker="o", markersize=4.8, label=MC_LABELS[note_type])
        if series.notna().any():
            last_idx = np.where(series.notna().to_numpy())[0][-1]
            ax.text(last_idx + 0.15, series.iloc[last_idx], fmt_pct(series.iloc[last_idx]), color=TYPE_COLORS[note_type], fontsize=8.5, fontweight="bold", va="center")
    ax.set_title("Quarterly share of documents with service opinions", loc="left", color=COLORS["ink"], fontsize=12, fontweight="bold")
    ax.set_ylabel("% with at least one service opinion", color=COLORS["muted"], fontsize=9)
    ax.set_ylim(0, max(100, np.nanmax(grouped["share"]) * 1.08))
    tick_step = max(1, math.ceil(len(quarters) / 10))
    ax.set_xticks(x[::tick_step], quarters[::tick_step])
    clean_axis(ax)
    ax.legend(ncols=3, loc="upper right", frameon=False, fontsize=9)
    return save_fig(fig, "other_quarterly_opinion_coverage")


def mc_note_detail_rows(df: pd.DataFrame) -> list[dict]:
    rows = []
    for note_type in MC_NOTE_TYPES:
        items = valid_mc_df(df)
        items = items[items["MC_Note_Type"].eq(note_type)].copy()
        if items.empty:
            continue
        items["Has Service Opinion"] = items.apply(has_service_opinion, axis=1)
        mean_total = items["Document Page Count"].mean()
        mean_pre = items["Page count before opinion"].mean()
        mean_annex = items["Annex Page Count"].mean()
        rows.append(
            {
                "type": MC_LABELS[note_type],
                "documents": len(items),
                "opinion_docs": int(items["Has Service Opinion"].sum()),
                "opinion_share": items["Has Service Opinion"].mean() * 100,
                "median_pages": items["Document Page Count"].median(),
                "pre_pages": mean_pre,
                "opinion_pages": max(0, mean_total - mean_pre - mean_annex),
                "annex_pages": mean_annex,
                "date_range": f"{items['BO Validation Date'].min().date()} to {items['BO Validation Date'].max().date()}",
            }
        )
    return rows


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(243, 246, 250)


def add_textbox(slide, x, y, w, h, text, font_size=18, color=COLORS["ink"], bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return box


def hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_title(slide, title: str, kicker: str = "OTHER focus"):
    add_textbox(slide, 0.48, 0.28, 2.2, 0.25, kicker.upper(), font_size=8.5, color=COLORS["muted"], bold=True)
    add_textbox(slide, 0.48, 0.5, 11.7, 0.45, title, font_size=24, color=COLORS["ink"], bold=True)
    line = slide.shapes.add_shape(1, Inches(0.48), Inches(1.02), Inches(12.35), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(COLORS["yellow"])
    line.line.color.rgb = hex_to_rgb(COLORS["yellow"])


def add_footer(slide):
    add_textbox(slide, 0.48, 7.18, 6.0, 0.18, "Source: Data Analysis/dashboard_data.json | Timeline: BO validation date", font_size=7.5, color=COLORS["muted"])


def add_callout(slide, text: str, x=9.65, y=1.34, w=2.8, h=4.6):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(214, 222, 232)
    tf = shape.text_frame
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.14)
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgb(COLORS["ink"])


def add_chart_slide(prs, title: str, image_path: Path, callout: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, title)
    slide.shapes.add_picture(str(image_path), Inches(0.52), Inches(1.2), width=Inches(9.0))
    add_callout(slide, callout)
    add_footer(slide)
    return slide


def add_kpi_card(slide, x, y, w, label, value, note=None):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(1.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(214, 222, 232)
    add_textbox(slide, x + 0.12, y + 0.12, w - 0.24, 0.18, label.upper(), font_size=7.8, color=COLORS["muted"], bold=True)
    add_textbox(slide, x + 0.12, y + 0.34, w - 0.24, 0.32, value, font_size=19, color=COLORS["ink"], bold=True)
    if note:
        add_textbox(slide, x + 0.12, y + 0.72, w - 0.24, 0.18, note, font_size=7.8, color=COLORS["muted"])


def add_title_slide(prs, df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_textbox(slide, 0.62, 0.58, 1.8, 0.25, "OTHER FOCUS", font_size=9, color=COLORS["muted"], bold=True)
    add_textbox(slide, 0.62, 0.9, 9.9, 0.7, "MC Notes story pack", font_size=34, color=COLORS["ink"], bold=True)
    add_textbox(slide, 0.62, 1.55, 9.5, 0.35, "A first presentation cut built from the dashboard data, centered on OTHER notes.", font_size=14, color=COLORS["muted"])

    opinions = service_opinion_frame(df)
    docs_with_opinions = opinions["file"].nunique()
    classified = valid_mc_df(df)
    dates = df["BO Validation Date"].dropna()
    add_kpi_card(slide, 0.62, 2.42, 2.35, "Documents", fmt_int(len(df)), "Template Type = OTHER")
    add_kpi_card(slide, 3.15, 2.42, 2.35, "Opinion rows", fmt_int(len(opinions)), "Across service columns")
    add_kpi_card(slide, 5.68, 2.42, 2.35, "Docs with opinions", fmt_int(docs_with_opinions), f"{fmt_pct(docs_with_opinions / len(df) * 100)} of OTHER")
    add_kpi_card(slide, 8.21, 2.42, 2.35, "Classified MC type", fmt_int(len(classified)), f"{fmt_pct(len(classified) / len(df) * 100)} of OTHER")
    add_kpi_card(slide, 10.74, 2.42, 2.35, "BO date range", f"{dates.min().date()}",
                 f"to {dates.max().date()}")

    add_textbox(slide, 0.62, 4.06, 11.8, 0.38, "Storyline", font_size=18, color=COLORS["ink"], bold=True)
    story = (
        "1. Service opinions are concentrated in a handful of services, but length patterns differ.\n"
        "2. Monthly page load is mostly driven by annexes, so page composition matters more than total pages alone.\n"
        "3. Decision notes dominate OTHER volume; opinion coverage is much lower for Info and Discussion notes.\n"
        "4. The deck uses BO validation dates for all time series."
    )
    add_textbox(slide, 0.72, 4.56, 11.5, 1.45, story, font_size=14, color=COLORS["ink"])
    add_footer(slide)


def add_detail_slide(prs, df: pd.DataFrame):
    rows = mc_note_detail_rows(df)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "MC Note Type Detail")

    headers = ["Type", "Docs", "Opinion docs", "Opinion share", "Median pages", "Pre-op", "Opinion", "Annex", "BO date range"]
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.55), Inches(1.28), Inches(12.1), Inches(2.15)).table
    widths = [1.15, 0.75, 1.05, 1.1, 1.0, 0.78, 0.78, 0.78, 3.0]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(232, 239, 246)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(8.2)
            p.font.color.rgb = hex_to_rgb(COLORS["ink"])
    for row_idx, item in enumerate(rows, start=1):
        values = [
            item["type"],
            fmt_int(item["documents"]),
            fmt_int(item["opinion_docs"]),
            fmt_pct(item["opinion_share"]),
            fmt_one(item["median_pages"]),
            fmt_one(item["pre_pages"]),
            fmt_one(item["opinion_pages"]),
            fmt_one(item["annex_pages"]),
            item["date_range"],
        ]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8.4)
                p.font.color.rgb = hex_to_rgb(COLORS["ink"])
                if col in range(1, 8):
                    p.alignment = PP_ALIGN.RIGHT

    ranked = sorted(rows, key=lambda item: item["opinion_share"], reverse=True)
    insight = (
        f"{ranked[0]['type']} has the highest opinion coverage at {fmt_pct(ranked[0]['opinion_share'])}, "
        f"while {ranked[-1]['type']} has the lowest at {fmt_pct(ranked[-1]['opinion_share'])}. "
        "That makes note type a useful first lens for explaining where service opinions appear in OTHER."
    )
    add_callout(slide, insight, x=0.68, y=4.1, w=11.6, h=1.0)
    add_footer(slide)


def main():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df = load_other_rows()

    service_chart = build_service_volume_chart(df)
    page_chart = build_page_composition_chart(df)
    mc_month_chart = build_mc_monthly_chart(df)
    quarterly_chart = build_quarterly_opinion_chart(df)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, df)

    opinions = service_opinion_frame(df)
    service_counts = Counter(opinions["service"])
    top_service, top_count = service_counts.most_common(1)[0]
    callout = (
        f"{top_service} is the highest-volume service with {fmt_int(top_count)} opinion rows. "
        "The right-hand markers separate volume from length, so services with fewer opinions but longer text still stand out."
    )
    add_chart_slide(prs, "Service Volume and Length", service_chart, callout)

    page_stats = page_composition_monthly(df)
    total_mean = df["Document Page Count"].mean()
    annex_mean = df["Annex Page Count"].mean()
    callout = (
        f"OTHER averages {fmt_one(total_mean)} total pages per document, with annexes averaging {fmt_one(annex_mean)} pages. "
        "The monthly view shows whether document growth is coming from pre-opinion text, opinions, or annex load."
    )
    add_chart_slide(prs, "Monthly Page Composition", page_chart, callout)

    classified = valid_mc_df(df)
    type_counts = classified["MC_Note_Type"].value_counts()
    dominant_type = type_counts.index[0]
    callout = (
        f"{MC_LABELS[dominant_type]} notes are the dominant classified OTHER note type "
        f"({fmt_int(type_counts.iloc[0])} documents). The moving-average lines make changes in the smaller types easier to see."
    )
    add_chart_slide(prs, "Monthly Documents By MC Note Type", mc_month_chart, callout)

    detail = mc_note_detail_rows(df)
    decision = next(item for item in detail if item["type"] == "Decision")
    info = next(item for item in detail if item["type"] == "Info")
    callout = (
        f"Opinion coverage differs sharply by type: Decision is {fmt_pct(decision['opinion_share'])}, "
        f"while Info is {fmt_pct(info['opinion_share'])}. This slide frames coverage as a rate, not a raw count."
    )
    add_chart_slide(prs, "Quarterly Opinion Coverage", quarterly_chart, callout)

    add_detail_slide(prs, df)

    prs.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")
    print(f"Assets saved in {ASSET_DIR}")


if __name__ == "__main__":
    main()
